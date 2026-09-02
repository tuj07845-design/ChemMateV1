from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path
from typing import Any

CN_SLIDE_FONT = "微软雅黑"     # PPT 全文

# ============================================================
# PPT 渲染
# ============================================================

# glm —— PPT run 中西文字体（a:latin + a:ea 同设微软雅黑）
def _pptx_run_font(run, size=None, bold=None, color=None, font: str = CN_SLIDE_FONT) -> None:
    """给 PPT 文字 run 设置中西文字体（a:latin + a:ea）、字号、加粗、颜色。"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    run.font.name = font  # a:latin
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        latin = rpr.find(qn("a:latin"))
        if latin is not None:
            latin.addnext(ea)
        else:
            rpr.append(ea)
    ea.set("typeface", font)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _render_pptx(title: str, sections: list, output_path: Path, author: str) -> Path:
    """渲染成 .pptx。每个 heading 开新页，同页内累积后续内容。
    排版约定见模块文档：全文微软雅黑，正文 18pt 行距 1.2，表头加粗。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # glm —— 16:9 页面
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # glm —— 封面：主标题 40pt 加粗居中 + 生成时间/作者
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    _pptx_run_font(r, size=40, bold=True)
    # 副信息
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}    {author}"
    _pptx_run_font(r2, size=14, color=RGBColor(0x66, 0x66, 0x66))

    # glm —— 内容页：每个 heading 开新页；内容溢出自动开续页
    current = None  # 当前 slide

    def new_slide(heading_text: str = ""):
        nonlocal current, cur_heading
        current = prs.slides.add_slide(blank)
        # glm —— 记录当前页标题（续页叠加时先去掉已有的“（续）”）
        cur_heading = heading_text[:-3] if heading_text.endswith("（续）") else heading_text
        # 顶部标题
        if heading_text:
            tbh = current.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
            tfh = tbh.text_frame
            tfh.word_wrap = True
            ph = tfh.paragraphs[0]
            rh = ph.add_run()
            rh.text = heading_text
            _pptx_run_font(rh, size=28, bold=True)

    # 内容区起始 y（标题下方）
    body_top = Inches(1.4)
    body_left = Inches(0.6)
    body_width = Inches(12.1)
    cur_y = body_top

    def reset_body():
        nonlocal cur_y
        cur_y = body_top

    # glm —— 新增：溢出守卫——内容放不下时自动开续页，避免超出 7.5 英寸画面被裁掉
    slide_bottom = Inches(7.1)
    cur_heading = ""

    def ensure_space(needed) -> None:
        nonlocal cur_y
        if int(cur_y) + int(needed) > int(slide_bottom):
            new_slide((cur_heading or "") + "（续）")
            reset_body()

    for sec in sections:
        stype = sec.get("type")

        if stype == "heading":
            # 开新页
            new_slide(str(sec.get("text", "")))
            reset_body()

        else:
            if current is None:
                new_slide("")
                reset_body()

            if stype == "paragraph":
                # glm —— 溢出守卫 + 正文 18pt 行距 1.2
                ensure_space(Inches(0.8))
                tbx = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(0.8))
                tfx = tbx.text_frame
                tfx.word_wrap = True
                pp = tfx.paragraphs[0]
                pp.line_spacing = 1.2
                rr = pp.add_run()
                rr.text = str(sec.get("text", ""))
                _pptx_run_font(rr, size=18)
                cur_y = Emu(int(cur_y) + Inches(0.6))

            elif stype == "bullets":
                items = sec.get("items") or []
                if isinstance(items, list):
                    # glm —— 溢出守卫（按条目数估高）
                    ensure_space(Inches(0.2 + 0.4 * max(1, len(items))))
                    tbx = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(1.2))
                    tfx = tbx.text_frame
                    tfx.word_wrap = True
                    first = True
                    for it in items:
                        para = tfx.paragraphs[0] if first else tfx.add_paragraph()
                        first = False
                        para.line_spacing = 1.2
                        para.space_after = Pt(6)
                        rr = para.add_run()
                        rr.text = "• " + str(it)
                        _pptx_run_font(rr, size=18)
                    cur_y = Emu(int(cur_y) + Inches(0.2) + len(items) * Inches(0.4))

            elif stype == "table":
                headers = sec.get("headers") or []
                rows = sec.get("rows") or []
                if not isinstance(headers, list) or not isinstance(rows, list):
                    continue
                n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
                if n_cols == 0:
                    continue
                n_rows = 1 + len(rows)
                # 估算表格高度
                tbl_h = Inches(0.4 + 0.35 * n_rows)
                # glm —— 溢出守卫：整表放不下就开续页
                ensure_space(int(tbl_h) + Inches(0.1))
                tbl_shape = current.shapes.add_table(
                    n_rows, n_cols, body_left, cur_y, body_width, tbl_h
                )
                table = tbl_shape.table
                # glm —— 表头 16pt 加粗、表体 14pt（微软雅黑）
                for j, htext in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = str(htext)
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            _pptx_run_font(r, size=16, bold=True)
                for i, row in enumerate(rows):
                    if not isinstance(row, list):
                        continue
                    for j in range(n_cols):
                        val = row[j] if j < len(row) else ""
                        table.cell(i + 1, j).text = "" if val is None else str(val)
                        for p in table.cell(i + 1, j).text_frame.paragraphs:
                            for r in p.runs:
                                _pptx_run_font(r, size=14)
                cur_y = Emu(int(cur_y) + int(tbl_h) + Inches(0.2))

            elif stype == "image":
                img_path = sec.get("path", "")
                caption = sec.get("caption", "")
                if not img_path or not os.path.isfile(img_path):
                    tbx = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(0.5))
                    tfx = tbx.text_frame
                    pp = tfx.paragraphs[0]
                    rr = pp.add_run()
                    rr.text = f"[图片缺失：{img_path}]"
                    _pptx_run_font(rr, size=14, color=RGBColor(0xCC, 0x00, 0x00))
                    cur_y = Emu(int(cur_y) + Inches(0.5))
                else:
                    # 图片高度限制，保持比例
                    try:
                        from PIL import Image
                        with Image.open(img_path) as im:
                            iw, ih = im.size
                    except Exception:
                        iw, ih = 900, 420
                    max_h = Inches(4.6)
                    max_w = body_width
                    # 按比例缩放
                    ratio = min(max_w / Inches(iw / 96.0), max_h / Inches(ih / 96.0)) if iw else 1
                    w = Inches(iw / 96.0 * ratio) if iw else max_w
                    h = Inches(ih / 96.0 * ratio) if ih else max_h
                    # glm —— 溢出守卫：图 + 图注放不下就开续页
                    ensure_space(int(h) + Inches(0.5))
                    try:
                        current.shapes.add_picture(img_path, body_left, cur_y, width=w, height=h)
                        cur_y = Emu(int(cur_y) + int(h) + Inches(0.1))
                    except Exception:
                        current.shapes.add_picture(img_path, body_left, cur_y, width=max_w)
                        cur_y = Emu(int(cur_y) + int(max_h) + Inches(0.1))
                    if caption:
                        cb = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(0.4))
                        cfp = cb.text_frame.paragraphs[0]
                        cfp.alignment = PP_ALIGN.CENTER
                        cr = cfp.add_run()
                        cr.text = str(caption)
                        _pptx_run_font(cr, size=12, color=RGBColor(0x66, 0x66, 0x66))
                        cur_y = Emu(int(cur_y) + Inches(0.4))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path

