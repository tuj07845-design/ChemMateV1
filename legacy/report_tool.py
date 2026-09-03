# -*- coding: utf-8 -*-
"""
ChemMate V1 — 报告生成工具（Word / PPT）

Agent 闭环最后一步：把 data_get / analyze / draw_mat 的结果
组织成 sections，由本工具渲染成 .docx 或 .pptx 文件。

分层定位（与 draw_mat / analyze_process 一致）：
    - 本工具只做"按 sections 渲染成文件"（数据工程），
      不跑 Aspen、不做分析、不调 MATLAB。
    - 报告的文字内容、结论由 Agent（LLM）组织后传入。
    - 图片由 draw_mat 生成，Agent 把 image_path 放进 image section。

用法（Agent 调用）：
    report_create(
        report_type="docx",            # 或 "pptx"
        title="10万吨环己烷 流程分析报告",
        sections=[
            {"type": "heading", "level": 1, "text": "一、流程概述"},
            {"type": "paragraph", "text": "甲苯加氢制环己烷……"},
            {"type": "bullets", "items": ["进料：甲苯+氢气", "反应器：B5"]},
            {"type": "table", "headers": ["流股","T/℃","P/bar"],
             "rows": [["S5",250,18],["S8",30,8]]},
            {"type": "image", "path": "jobs/.../figure.png", "caption": "流股 T/P"},
        ],
        output_path=None,              # 默认放当前目录下 reports/
    )

返回：
    {"success": True, "file_path": "...", "report_type": "docx",
     "section_count": 5, "error": "", "message": ""}

排版约定（符合中文阅读习惯）：
    - Word：A4 页面；正文宋体小四 + Times New Roman，1.5 倍行距，首行缩进两字符；
      标题黑体加粗（二号/三号/四号/小四）；表格宋体五号，表头加粗居中、跨页重复；
      页脚“第 X 页”。
    - PPT：全文微软雅黑；页标题加粗，正文 18pt、行距 1.2，表头 16pt 加粗；
      内容超出页面时自动开续页（标题带“（续）”）。

代码中带 “# glm” 注释的位置，即为按中文阅读习惯改造或新增的排版点。
"""

from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path
from typing import Any

TOOL_NAME = "report_create"

TOOL_DESCRIPTION = (
    "生成 Word(.docx) 或 PPT(.pptx) 分析报告。"
    "把要写进报告的内容按 sections 传入；"
    "图片用 image 类型，path 填 draw_mat 返回的 image_path。"
    "report_type 为 docx 或 pptx。"
    "输出已按中文阅读排版（Word：宋体正文/黑体标题/1.5 倍行距；PPT：微软雅黑）。"
)

TOOL_PARAMETERS = {
    "type": "object",
    "properties": {
        "report_type": {
            "type": "string",
            "enum": ["docx", "pptx"],
            "description": "报告格式：docx（Word）或 pptx（PowerPoint）",
        },
        "title": {
            "type": "string",
            "description": "报告标题",
        },
        "sections": {
            "type": "array",
            "description": "内容块列表，每块一个 dict，见模块文档",
            "items": {
                "type": "object",
                "properties": {
                    "type": {
                        "type": "string",
                        "enum": ["heading", "paragraph", "bullets", "table", "image"],
                    },
                    "level": {"type": "integer", "description": "heading 用，1~3"},
                    "text": {"type": "string", "description": "heading/paragraph 用"},
                    "items": {"type": "array", "items": {"type": "string"}, "description": "bullets 用"},
                    "headers": {"type": "array", "items": {"type": "string"}, "description": "table 用"},
                    "rows": {"type": "array", "description": "table 用，每行一个值列表"},
                    "path": {"type": "string", "description": "image 用，图片路径"},
                    "caption": {"type": "string", "description": "image 用，图注"},
                },
                "required": ["type"],
            },
        },
        "output_path": {
            "type": "string",
            "description": "输出文件路径；不传则默认放 reports/ 下",
        },
        "author": {
            "type": "string",
            "description": "作者/署名，默认 ChemMate V1",
        },
    },
    "required": ["report_type", "title", "sections"],
}

# 合法的 section 类型
SECTION_TYPES = ("heading", "paragraph", "bullets", "table", "image")


class ReportError(Exception):
    """报告生成失败。code 供上层稳定匹配。"""

    def __init__(self, code: str, message: str):
        super().__init__(message)
        self.code = code
        self.message = message


# ============================================================
# 缺依赖检测
# ============================================================

def _check_lib(report_type: str) -> None:
    """按格式检测依赖库，缺失给清晰指引。"""
    if report_type == "docx":
        try:
            import docx  # noqa: F401
        except ImportError:
            raise ReportError(
                "missing_dependency",
                "生成 Word 报告需要 python-docx，请执行：pip install python-docx",
            )
    elif report_type == "pptx":
        try:
            import pptx  # noqa: F401
        except ImportError:
            raise ReportError(
                "missing_dependency",
                "生成 PPT 报告需要 python-pptx，请执行：pip install python-pptx",
            )


# ============================================================
# section 校验
# ============================================================

def _validate_sections(sections: list) -> list:
    """校验并规整 sections；非法块抛 ReportError。"""
    if not isinstance(sections, list):
        raise ReportError("spec_invalid", "sections 必须是列表")

    cleaned = []
    for i, sec in enumerate(sections):
        if not isinstance(sec, dict):
            raise ReportError("spec_invalid", f"第 {i + 1} 块不是 dict")
        stype = sec.get("type")
        if stype not in SECTION_TYPES:
            raise ReportError(
                "spec_invalid",
                f"第 {i + 1} 块 type={stype!r} 非法，可选 {SECTION_TYPES}",
            )
        cleaned.append(sec)
    return cleaned


# ============================================================
# Word 渲染
# ============================================================

# glm —— 中西文字体集中定义：正文宋体/Times，标题黑体/Arial，PPT 微软雅黑
CN_BODY_FONT = "宋体"          # Word 正文中文
CN_HEAD_FONT = "黑体"          # Word 标题中文
CN_SLIDE_FONT = "微软雅黑"     # PPT 全文
LATIN_BODY_FONT = "Times New Roman"
LATIN_HEAD_FONT = "Arial"


# glm —— 段落样式级中西文字体（w:rFonts/@w:eastAsia）
def _docx_style_font(style, east: str = CN_BODY_FONT, latin: str = LATIN_BODY_FONT) -> None:
    """给段落样式设置中西文字体（w:rFonts/@w:eastAsia 决定中文字体）。"""
    from docx.oxml.ns import qn
    style.font.name = latin  # 先设西文，确保 rFonts 存在
    style.element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), east)


# glm —— run 级中西文字体/字号/加粗/颜色
def _docx_run_font(run, size=None, bold=None, color=None,
                   east: str = CN_BODY_FONT, latin: str = LATIN_BODY_FONT) -> None:
    """给 run 设置中西文字体、字号、加粗、颜色。"""
    from docx.oxml.ns import qn
    run.font.name = latin
    run._element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), east)
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


# glm —— 首行缩进按字符数（firstLineChars 为主，中文两字符习惯）
def _docx_first_line_indent(paragraph, chars: int = 2) -> None:
    """首行缩进按字符数：firstLineChars 为准（中文 Word 习惯），firstLine 兜底。"""
    from docx.oxml.ns import qn
    ind = paragraph._p.get_or_add_pPr().get_or_add_ind()
    ind.set(qn("w:firstLine"), str(chars * 240))       # 12pt 字号时一字符 240 twips
    ind.set(qn("w:firstLineChars"), str(chars * 100))


# glm —— 页脚居中“第 X 页”（PAGE 域）
def _docx_add_page_number(footer_paragraph) -> None:
    """页脚居中插入“第 X 页”（PAGE 域）。"""
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    footer_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _docx_run_font(footer_paragraph.add_run("第 "), size=Pt(9))
    r = footer_paragraph.add_run()
    fld_b = OxmlElement("w:fldChar")
    fld_b.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "
    fld_e = OxmlElement("w:fldChar")
    fld_e.set(qn("w:fldCharType"), "end")
    r._r.append(fld_b)
    r._r.append(instr)
    r._r.append(fld_e)
    r.font.size = Pt(9)
    _docx_run_font(footer_paragraph.add_run(" 页"), size=Pt(9))


def _render_docx(title: str, sections: list, output_path: Path, author: str) -> Path:
    """渲染成 .docx（排版约定见模块文档：宋体正文/黑体标题/1.5 倍行距/首行缩进）。"""
    from docx import Document
    from docx.shared import Pt, Inches, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import OxmlElement

    doc = Document()

    # glm —— 页面：A4，上下 2.54cm，左右 3.18cm（Word 常规页边距）
    sec0 = doc.sections[0]
    sec0.page_width, sec0.page_height = Cm(21.0), Cm(29.7)
    sec0.top_margin = sec0.bottom_margin = Cm(2.54)
    sec0.left_margin = sec0.right_margin = Cm(3.18)

    # glm —— 样式：正文宋体小四、1.5 倍行距；标题黑体加粗；项目符号宋体
    normal = doc.styles["Normal"]
    _docx_style_font(normal, CN_BODY_FONT, LATIN_BODY_FONT)
    normal.font.size = Pt(12)
    normal.paragraph_format.line_spacing = 1.5
    normal.paragraph_format.space_after = Pt(0)

    for style_name, size, before, after in (
        ("Title", 22, 0, 12),
        ("Heading 1", 16, 14, 8),
        ("Heading 2", 14, 10, 6),
        ("Heading 3", 12, 8, 4),
    ):
        st = doc.styles[style_name]
        _docx_style_font(st, CN_HEAD_FONT, LATIN_HEAD_FONT)
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = RGBColor(0, 0, 0)
        st.paragraph_format.line_spacing = 1.5
        st.paragraph_format.space_before = Pt(before)
        st.paragraph_format.space_after = Pt(after)
        st.paragraph_format.keep_with_next = True

    list_bullet = doc.styles["List Bullet"]
    _docx_style_font(list_bullet, CN_BODY_FONT, LATIN_BODY_FONT)
    list_bullet.font.size = Pt(12)
    list_bullet.paragraph_format.line_spacing = 1.5
    list_bullet.paragraph_format.space_after = Pt(0)

    # 核心属性
    try:
        doc.core_properties.author = author
        doc.core_properties.title = title
    except Exception:
        pass

    # glm —— 封面：标题居中（Title 样式）
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # glm —— 副信息：生成时间 + 作者（小字灰色居中）
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _docx_run_font(
        meta.add_run(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}    {author}"),
        size=Pt(10), color=RGBColor(0x66, 0x66, 0x66),
    )

    # glm —— 页脚页码
    _docx_add_page_number(sec0.footer.paragraphs[0])

    for sec in sections:
        stype = sec.get("type")

        if stype == "heading":
            level = int(sec.get("level", 1) or 1)
            level = max(1, min(3, level))
            doc.add_heading(str(sec.get("text", "")), level=level)

        elif stype == "paragraph":
            # glm —— 正文：两端对齐 + 首行缩进两字符
            p = doc.add_paragraph(str(sec.get("text", "")))
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            _docx_first_line_indent(p, 2)

        elif stype == "bullets":
            items = sec.get("items") or []
            if isinstance(items, list):
                for it in items:
                    bp = doc.add_paragraph(str(it), style="List Bullet")
                    bp.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        elif stype == "table":
            headers = sec.get("headers") or []
            rows = sec.get("rows") or []
            if not isinstance(headers, list) or not isinstance(rows, list):
                continue
            n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
            if n_cols == 0:
                continue
            table = doc.add_table(rows=1 + len(rows), cols=n_cols)
            table.style = "Light Grid Accent 1"
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            # glm —— 表格排版：宋体五号，表头加粗居中；行禁止跨页断开
            for i, row in enumerate([headers] + list(rows)):
                vals = row if isinstance(row, list) else []
                trpr = table.rows[i]._tr.get_or_add_trPr()
                trpr.append(OxmlElement("w:cantSplit"))
                for j in range(n_cols):
                    val = vals[j] if j < len(vals) else ""
                    cell = table.rows[i].cells[j]
                    cell.text = "" if val is None else str(val)
                    for p in cell.paragraphs:
                        p.alignment = (WD_ALIGN_PARAGRAPH.CENTER if i == 0
                                       else WD_ALIGN_PARAGRAPH.LEFT)
                        p.paragraph_format.line_spacing = 1.0
                        for r in p.runs:
                            _docx_run_font(r, size=Pt(10.5), bold=(i == 0))
            # glm —— 表头跨页重复
            table.rows[0]._tr.get_or_add_trPr().append(OxmlElement("w:tblHeader"))

        elif stype == "image":
            img_path = sec.get("path", "")
            caption = sec.get("caption", "")
            if not img_path or not os.path.isfile(img_path):
                # 图片缺失不中断，插一行提示
                p = doc.add_paragraph()
                _docx_run_font(p.add_run(f"[图片缺失：{img_path}]"),
                               color=RGBColor(0xCC, 0x00, 0x00))
            else:
                try:
                    doc.add_picture(img_path, width=Inches(5.5))
                except Exception:
                    doc.add_picture(img_path)
                # glm —— 图注：小字灰色居中
                if caption:
                    cap = doc.add_paragraph(str(caption))
                    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for r in cap.runs:
                        _docx_run_font(r, size=Pt(9), color=RGBColor(0x66, 0x66, 0x66))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    return output_path


# ============================================================
# PPT 渲染
# ============================================================

# glm —— PPT run 中西文字体（a:latin + a:ea 同设微软雅黑）
def _pptx_run_font(run, size=None, bold=None, color=None, font: str = CN_SLIDE_FONT) -> None:
    """给 PPT 文字 run 设置中西文字体（a:latin + a:ea）、字号、加粗、颜色。"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    run.font.name = font  # a:latin
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        latin = rpr.find(qn("a:latin"))
        if latin is not None:
            latin.addnext(ea)
        else:
            rpr.append(ea)
    ea.set("typeface", font)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _render_pptx(title: str, sections: list, output_path: Path, author: str) -> Path:
    """渲染成 .pptx。每个 heading 开新页，同页内累积后续内容。
    排版约定见模块文档：全文微软雅黑，正文 18pt 行距 1.2，表头加粗。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # glm —— 16:9 页面
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # glm —— 封面：主标题 40pt 加粗居中 + 生成时间/作者
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    _pptx_run_font(r, size=40, bold=True)
    # 副信息
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}    {author}"
    _pptx_run_font(r2, size=14, color=RGBColor(0x66, 0x66, 0x66))

    # glm —— 内容页：每个 heading 开新页；内容溢出自动开续页
    current = None  # 当前 slide

    def new_slide(heading_text: str = ""):
        nonlocal current, cur_heading
        current = prs.slides.add_slide(blank)
        # glm —— 记录当前页标题（续页叠加时先去掉已有的“（续）”）
        cur_heading = heading_text[:-3] if heading_text.endswith("（续）") else heading_text
        # 顶部标题
        if heading_text:
            tbh = current.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
            tfh = tbh.text_frame
            tfh.word_wrap = True
            ph = tfh.paragraphs[0]
            rh = ph.add_run()
            rh.text = heading_text
            _pptx_run_font(rh, size=28, bold=True)

    # 内容区起始 y（标题下方）
    body_top = Inches(1.4)
    body_left = Inches(0.6)
    body_width = Inches(12.1)
    cur_y = body_top

    def reset_body():
        nonlocal cur_y
        cur_y = body_top

    # glm —— 新增：溢出守卫——内容放不下时自动开续页，避免超出 7.5 英寸画面被裁掉
    slide_bottom = Inches(7.1)
    cur_heading = ""

    def ensure_space(needed) -> None:
        nonlocal cur_y
        if int(cur_y) + int(needed) > int(slide_bottom):
            new_slide((cur_heading or "") + "（续）")
            reset_body()

    for sec in sections:
        stype = sec.get("type")

        if stype == "heading":
            # 开新页
            new_slide(str(sec.get("text", "")))
            reset_body()

        else:
            if current is None:
                new_slide("")
                reset_body()

            if stype == "paragraph":
                # glm —— 溢出守卫 + 正文 18pt 行距 1.2
                ensure_space(Inches(0.8))
                tbx = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(0.8))
                tfx = tbx.text_frame
                tfx.word_wrap = True
                pp = tfx.paragraphs[0]
                pp.line_spacing = 1.2
                rr = pp.add_run()
                rr.text = str(sec.get("text", ""))
                _pptx_run_font(rr, size=18)
                cur_y = Emu(int(cur_y) + Inches(0.6))

            elif stype == "bullets":
                items = sec.get("items") or []
                if isinstance(items, list):
                    # glm —— 溢出守卫（按条目数估高）
                    ensure_space(Inches(0.2 + 0.4 * max(1, len(items))))
                    tbx = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(1.2))
                    tfx = tbx.text_frame
                    tfx.word_wrap = True
                    first = True
                    for it in items:
                        para = tfx.paragraphs[0] if first else tfx.add_paragraph()
                        first = False
                        para.line_spacing = 1.2
                        para.space_after = Pt(6)
                        rr = para.add_run()
                        rr.text = "• " + str(it)
                        _pptx_run_font(rr, size=18)
                    cur_y = Emu(int(cur_y) + Inches(0.2) + len(items) * Inches(0.4))

            elif stype == "table":
                headers = sec.get("headers") or []
                rows = sec.get("rows") or []
                if not isinstance(headers, list) or not isinstance(rows, list):
                    continue
                n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
                if n_cols == 0:
                    continue
                n_rows = 1 + len(rows)
                # 估算表格高度
                tbl_h = Inches(0.4 + 0.35 * n_rows)
                # glm —— 溢出守卫：整表放不下就开续页
                ensure_space(int(tbl_h) + Inches(0.1))
                tbl_shape = current.shapes.add_table(
                    n_rows, n_cols, body_left, cur_y, body_width, tbl_h
                )
                table = tbl_shape.table
                # glm —— 表头 16pt 加粗、表体 14pt（微软雅黑）
                for j, htext in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = str(htext)
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            _pptx_run_font(r, size=16, bold=True)
                for i, row in enumerate(rows):
                    if not isinstance(row, list):
                        continue
                    for j in range(n_cols):
                        val = row[j] if j < len(row) else ""
                        table.cell(i + 1, j).text = "" if val is None else str(val)
                        for p in table.cell(i + 1, j).text_frame.paragraphs:
                            for r in p.runs:
                                _pptx_run_font(r, size=14)
                cur_y = Emu(int(cur_y) + int(tbl_h) + Inches(0.2))

            elif stype == "image":
                img_path = sec.get("path", "")
                caption = sec.get("caption", "")
                if not img_path or not os.path.isfile(img_path):
                    tbx = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(0.5))
                    tfx = tbx.text_frame
                    pp = tfx.paragraphs[0]
                    rr = pp.add_run()
                    rr.text = f"[图片缺失：{img_path}]"
                    _pptx_run_font(rr, size=14, color=RGBColor(0xCC, 0x00, 0x00))
                    cur_y = Emu(int(cur_y) + Inches(0.5))
                else:
                    # 图片高度限制，保持比例
                    try:
                        from PIL import Image
                        with Image.open(img_path) as im:
                            iw, ih = im.size
                    except Exception:
                        iw, ih = 900, 420
                    max_h = Inches(4.6)
                    max_w = body_width
                    # 按比例缩放
                    ratio = min(max_w / Inches(iw / 96.0), max_h / Inches(ih / 96.0)) if iw else 1
                    w = Inches(iw / 96.0 * ratio) if iw else max_w
                    h = Inches(ih / 96.0 * ratio) if ih else max_h
                    # glm —— 溢出守卫：图 + 图注放不下就开续页
                    ensure_space(int(h) + Inches(0.5))
                    try:
                        current.shapes.add_picture(img_path, body_left, cur_y, width=w, height=h)
                        cur_y = Emu(int(cur_y) + int(h) + Inches(0.1))
                    except Exception:
                        current.shapes.add_picture(img_path, body_left, cur_y, width=max_w)
                        cur_y = Emu(int(cur_y) + int(max_h) + Inches(0.1))
                    if caption:
                        cb = current.shapes.add_textbox(body_left, cur_y, body_width, Inches(0.4))
                        cfp = cb.text_frame.paragraphs[0]
                        cfp.alignment = PP_ALIGN.CENTER
                        cr = cfp.add_run()
                        cr.text = str(caption)
                        _pptx_run_font(cr, size=12, color=RGBColor(0x66, 0x66, 0x66))
                        cur_y = Emu(int(cur_y) + Inches(0.4))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# ============================================================
# 输出路径
# ============================================================

def _default_output_path(report_type: str, title: str) -> Path:
    """默认输出到 cwd/reports/ 下，文件名用时间戳。"""
    root = Path(os.environ.get("CHEMMATE_REPORTS_DIR") or Path.cwd() / "reports")
    root.mkdir(parents=True, exist_ok=True)
    # 文件名：标题前若干字 + 时间戳，去非法字符
    safe = "".join(c for c in title if c not in r'\/:*?"<>|')[:20].strip() or "report"
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return root / f"{safe}_{stamp}.{report_type}"


# ============================================================
# 主入口
# ============================================================

def report_create(
    report_type: str,
    title: str,
    sections: list,
    output_path: str | Path | None = None,
    author: str = "ChemMate V1",
) -> dict[str, Any]:
    """生成 Word / PPT 报告。

    返回统一 dict：
        success / error / message / file_path / report_type / section_count
    """
    rt = (report_type or "").strip().lower()
    if rt not in ("docx", "pptx"):
        return _fail("spec_invalid", "report_type 必须是 docx 或 pptx")

    if not title or not str(title).strip():
        return _fail("spec_invalid", "title 不能为空")

    try:
        sections = _validate_sections(sections)
    except ReportError as exc:
        return _fail(exc.code, exc.message)

    try:
        _check_lib(rt)
    except ReportError as exc:
        return _fail(exc.code, exc.message)

    out = Path(output_path) if output_path else _default_output_path(rt, str(title))
    # 保证扩展名
    if out.suffix.lower() != f".{rt}":
        out = out.with_suffix(f".{rt}")

    try:
        if rt == "docx":
            _render_docx(str(title), sections, out, author)
        else:
            _render_pptx(str(title), sections, out, author)
    except ReportError as exc:
        return _fail(exc.code, exc.message)
    except Exception as exc:
        return _fail("render_failed", f"{type(exc).__name__}: {exc}")

    return {
        "success": True,
        "error": "",
        "message": f"已生成 {rt} 报告，共 {len(sections)} 个内容块",
        "file_path": str(out),
        "report_type": rt,
        "section_count": len(sections),
    }


def _fail(code: str, message: str) -> dict:
    return {
        "success": False,
        "error": code,
        "message": message,
        "file_path": "",
        "report_type": "",
        "section_count": 0,
    }


# ============================================================
# 便捷：从 analyze 结果 + 图片快速组装 sections
# ============================================================

def build_sections_from_results(
    process_data: dict | None = None,
    analyze_result: dict | None = None,
    image_paths: list[str] | None = None,
    intro: str = "",
) -> list[dict]:
    """把 data_get / analyze / draw_mat 的结果快速拼成 sections。

    Agent 也可自己组织 sections，本函数只是省事用的便捷构造。
    返回的 sections 可直接传给 report_create。
    """
    sections: list[dict] = []

    if intro:
        sections.append({"type": "paragraph", "text": intro})

    # 流程概况表（来自 data_get）
    if isinstance(process_data, dict) and process_data.get("success"):
        streams = process_data.get("streams", {})
        if isinstance(streams, dict) and streams:
            sections.append({"type": "heading", "level": 1, "text": "一、流程概况"})
            blocks = process_data.get("blocks", [])
            conns = process_data.get("connections", [])
            sections.append({"type": "paragraph", "text": f"设备 {len(blocks)} 个，物流 {len(streams)} 条，连接 {len(conns)} 条。"})
            # 设备表
            if isinstance(conns, list) and conns:
                headers = ["设备", "输入", "输出"]
                rows = []
                for c in conns:
                    if not isinstance(c, dict):
                        continue
                    rows.append([
                        str(c.get("block", "")),
                        ", ".join(str(x) for x in (c.get("inputs") or [])),
                        ", ".join(str(x) for x in (c.get("outputs") or [])),
                    ])
                if rows:
                    sections.append({"type": "table", "headers": headers, "rows": rows})

    # 分析结果（来自 analyze_process）
    if isinstance(analyze_result, dict) and analyze_result.get("success"):
        sections.append({"type": "heading", "level": 1, "text": "二、数据分析"})

        summary = analyze_result.get("summary", {})
        if isinstance(summary, dict):
            items = []
            for k, v in summary.items():
                items.append(f"{k}：{v}")
            if items:
                sections.append({"type": "bullets", "items": items})

        findings = analyze_result.get("findings", [])
        if isinstance(findings, list) and findings:
            sections.append({"type": "heading", "level": 2, "text": "数据检查发现"})
            items = []
            for f in findings:
                if not isinstance(f, dict):
                    continue
                items.append(
                    f"[{f.get('level', '?')}] {f.get('type', '')}"
                    + (f"（{f.get('stream', '')}）" if f.get("stream") else "")
                    + f"：{f.get('message', '')}"
                )
            sections.append({"type": "bullets", "items": items})
        else:
            sections.append({"type": "paragraph", "text": "数据检查未发现异常。"})

        tracking = analyze_result.get("component_tracking", [])
        if isinstance(tracking, list) and tracking:
            sections.append({"type": "heading", "level": 2, "text": "组分追踪"})
            headers = ["流股", "摩尔分率", "摩尔流", "质量流"]
            rows = []
            for t in tracking:
                if not isinstance(t, dict):
                    continue
                rows.append([
                    str(t.get("stream", "")),
                    _fmt(t.get("mole_fraction")),
                    _fmt_unit(t.get("mole_flow"), t.get("mole_flow_unit")),
                    _fmt_unit(t.get("mass_flow"), t.get("mass_flow_unit")),
                ])
            sections.append({"type": "table", "headers": headers, "rows": rows})

        changes = analyze_result.get("stream_changes", [])
        if isinstance(changes, list) and changes:
            sections.append({"type": "heading", "level": 2, "text": "前后变化"})
            headers = ["设备", "从", "到", "组分变化"]
            rows = []
            for c in changes:
                if not isinstance(c, dict):
                    continue
                comps = c.get("significant_components", [])
                desc = "; ".join(
                    f"{x.get('component', '')}:{_fmt(x.get('from'))}→{_fmt(x.get('to'))}"
                    for x in comps if isinstance(x, dict)
                )
                rows.append([str(c.get("block", "")), str(c.get("from_stream", "")), str(c.get("to_stream", "")), desc])
            sections.append({"type": "table", "headers": headers, "rows": rows})

    # 图片
    if image_paths:
        sections.append({"type": "heading", "level": 1, "text": "三、图表"})
        for i, p in enumerate(image_paths, 1):
            sections.append({"type": "image", "path": str(p), "caption": f"图{i}"})

    return sections


def _fmt(v) -> str:
    if v is None:
        return ""
    try:
        return f"{float(v):.6g}"
    except (TypeError, ValueError):
        return str(v)


def _fmt_unit(v, u) -> str:
    s = _fmt(v)
    if s and u:
        return f"{s} {u}"
    return s


# ============================================================
# Agent Tool 注册信息
# ============================================================

def tool_spec() -> dict:
    return {
        "name": TOOL_NAME,
        "description": TOOL_DESCRIPTION,
        "parameters": TOOL_PARAMETERS,
    }


if __name__ == "__main__":
    # 自测：生成一个示例 docx
    r = report_create(
        report_type="docx",
        title="report_tool 自测",
        sections=[
            {"type": "heading", "level": 1, "text": "一、概述"},
            {"type": "paragraph", "text": "这是 report_tool 的自测报告。"},
            {"type": "bullets", "items": ["要点 A", "要点 B"]},
            {"type": "table", "headers": ["流股", "T/℃", "P/bar"], "rows": [["S5", 250, 18], ["S8", 30, 8]]},
        ],
    )
    print(r)
